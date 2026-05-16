from __future__ import annotations

import base64
import importlib
import json
import os
import sys
import tempfile
import unittest
from pathlib import Path
from typing import Any
from unittest import mock

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "src"))

from mcp.types import ImageContent, TextContent

from ntulearn_mcp import server
from ntulearn_mcp.client import BbRouterExpiredError, NTULearnClient

# Minimal hand-crafted PDF v1.4 with the visible text "Hello PDF Fixture".
# 598 bytes raw — kept as a base64 constant to avoid a fixture file and the
# chicken-and-egg of generating a fixture with the library under test.
_TINY_PDF_B64 = (
    "JVBERi0xLjQKJeLjz9MKMSAwIG9iago8PCAvVHlwZSAvQ2F0YWxvZyAvUGFnZXMgMiAwIFIgPj4K"
    "ZW5kb2JqCjIgMCBvYmoKPDwgL1R5cGUgL1BhZ2VzIC9LaWRzIFszIDAgUl0gL0NvdW50IDEgPj4K"
    "ZW5kb2JqCjMgMCBvYmoKPDwgL1R5cGUgL1BhZ2UgL1BhcmVudCAyIDAgUiAvTWVkaWFCb3ggWzAg"
    "MCA2MTIgNzkyXSAvQ29udGVudHMgNCAwIFIgL1Jlc291cmNlcyA8PCAvRm9udCA8PCAvRjEgNSAw"
    "IFIgPj4gPj4gPj4KZW5kb2JqCjQgMCBvYmoKPDwgL0xlbmd0aCA0OSA+PgpzdHJlYW0KQlQKL0Yx"
    "IDI0IFRmCjcyIDcyMCBUZAooSGVsbG8gUERGIEZpeHR1cmUpIFRqCkVUCmVuZHN0cmVhbQplbmRv"
    "YmoKNSAwIG9iago8PCAvVHlwZSAvRm9udCAvU3VidHlwZSAvVHlwZTEgL0Jhc2VGb250IC9IZWx2"
    "ZXRpY2EgPj4KZW5kb2JqCnhyZWYKMCA2CjAwMDAwMDAwMDAgNjU1MzUgZiAKMDAwMDAwMDAxNSAw"
    "MDAwMCBuIAowMDAwMDAwMDY0IDAwMDAwIG4gCjAwMDAwMDAxMjEgMDAwMDAgbiAKMDAwMDAwMDI0"
    "NyAwMDAwMCBuIAowMDAwMDAwMzQ1IDAwMDAwIG4gCnRyYWlsZXIKPDwgL1NpemUgNiAvUm9vdCAx"
    "IDAgUiA+PgpzdGFydHhyZWYKNDE1CiUlRU9GCg=="
)
_TINY_PDF_BYTES = base64.b64decode(_TINY_PDF_B64)


# --- Office-format fixture builders ----------------------------------------
# These generate fixtures at test time using the same library under test.
# Unlike PDF (where pypdf is the consumer), here the value is testing OUR
# extraction wrapper — round-tripping through the library is fine.

def _make_docx_bytes(
    paragraphs: list[str], table_rows: list[list[str]] | None = None
) -> bytes:
    from io import BytesIO

    from docx import Document

    doc = Document()
    for p in paragraphs:
        doc.add_paragraph(p)
    if table_rows:
        table = doc.add_table(rows=len(table_rows), cols=len(table_rows[0]))
        for i, row in enumerate(table_rows):
            for j, val in enumerate(row):
                table.rows[i].cells[j].text = val
    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_pptx_bytes(slides: list[dict[str, Any]]) -> bytes:
    """Build a deck. Each slide dict supports keys: title, body, notes."""
    from io import BytesIO

    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank = prs.slide_layouts[6]  # Blank layout — gives us full control
    for s in slides:
        slide = prs.slides.add_slide(blank)
        if "title" in s:
            tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
            tb.text_frame.text = s["title"]
        if "body" in s:
            tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(5), Inches(2))
            tb.text_frame.text = s["body"]
        if s.get("notes"):
            slide.notes_slide.notes_text_frame.text = s["notes"]
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_xlsx_bytes(sheets: dict[str, list[list[Any]]]) -> bytes:
    from io import BytesIO

    from openpyxl import Workbook

    wb = Workbook()
    # Remove the auto-created default sheet so test sheet names are exact.
    wb.remove(wb.active)
    for name, rows in sheets.items():
        ws = wb.create_sheet(name)
        for row in rows:
            ws.append(row)
    buf = BytesIO()
    wb.save(buf)
    return buf.getvalue()


class _CookieEnvIsolation:
    """Mixin that snapshots NTULEARN_COOKIE + server._client and restores them.

    Also stubs out browser auto-read and the keychain cache so tests never
    touch the user's real browser cookies or OS keychain, and never hit
    the retry/backoff sleeps inside ``read_bbrouter_cookie``. Tests that
    want to exercise specific browser/cache behavior re-patch the same
    symbols inside the test body — those override these defaults.
    """

    def setUp(self) -> None:
        self._old_env = os.environ.get("NTULEARN_COOKIE")
        self._old_client = server._client

        # Default: browser auto-read returns None (so resolution falls
        # through to env / cache deterministically) and cache is empty +
        # inert. Tests that need a different shape patch the same symbols
        # inside their `with` blocks.
        self._browser_patch = mock.patch.object(
            server, "read_bbrouter_cookie", return_value=None
        )
        self._cache_read_patch = mock.patch.object(
            server, "read_cached_cookie", return_value=None
        )
        self._cache_write_patch = mock.patch.object(
            server, "write_cached_cookie", return_value=False
        )
        self._cache_delete_patch = mock.patch.object(
            server, "delete_cached_cookie", return_value=None
        )
        self._browser_patch.start()
        self._cache_read_patch.start()
        self._cache_write_patch.start()
        self._cache_delete_patch.start()

    def tearDown(self) -> None:
        self._cache_delete_patch.stop()
        self._cache_write_patch.stop()
        self._cache_read_patch.stop()
        self._browser_patch.stop()
        if self._old_env is None:
            os.environ.pop("NTULEARN_COOKIE", None)
        else:
            os.environ["NTULEARN_COOKIE"] = self._old_env
        server._client = self._old_client


class MissingCookieTests(_CookieEnvIsolation, unittest.IsolatedAsyncioTestCase):
    async def test_missing_cookie_raises(self) -> None:
        # call_tool now raises on errors so the MCP framework can wrap them
        # into CallToolResult(isError=True). Verify the message is preserved.
        os.environ.pop("NTULEARN_COOKIE", None)
        server._client = None
        with mock.patch.object(server, "read_bbrouter_cookie", return_value=None):
            with self.assertRaises(RuntimeError) as ctx:
                await server.call_tool("ntulearn_list_courses", {})
        self.assertIn("No NTULearn cookie found", str(ctx.exception))


class DotenvPrecedenceTests(unittest.TestCase):
    def test_env_value_wins_over_dotenv(self) -> None:
        env_path = ROOT / ".env"
        old_env = os.environ.get("NTULEARN_COOKIE")
        old_env_file = env_path.read_text(encoding="utf-8") if env_path.exists() else None
        try:
            os.environ["NTULEARN_COOKIE"] = "fresh-from-env"
            env_path.write_text("NTULEARN_COOKIE=stale-from-dotenv\n", encoding="utf-8")
            reloaded = importlib.reload(server)
            # Resolution order is browser → env → cache → raise. To assert
            # env-from-environ beats env-from-dotenv we need to silence the
            # browser and cache paths; otherwise a real BbRouter cookie on
            # the dev machine (or a previously-cached value) would win.
            with mock.patch.object(reloaded, "read_bbrouter_cookie", return_value=None):
                with mock.patch.object(reloaded, "read_cached_cookie", return_value=None):
                    self.assertEqual(reloaded._resolve_cookie(), "fresh-from-env")
        finally:
            if old_env is None:
                os.environ.pop("NTULEARN_COOKIE", None)
            else:
                os.environ["NTULEARN_COOKIE"] = old_env

            if old_env_file is None:
                env_path.unlink(missing_ok=True)
            else:
                env_path.write_text(old_env_file, encoding="utf-8")

            importlib.reload(server)


class CookieResolutionTests(_CookieEnvIsolation, unittest.TestCase):
    """Resolution order is browser → env → cache → raise.

    Browser-first matches the MCP server's reason for existing: zero-config
    convenience when the user is logged in somewhere we can read. Env var
    is a manual fallback (e.g. Windows + Chrome/Edge ABE), cache is a
    keychain-backed safety net for transient browser failures.
    """

    def test_browser_takes_precedence_over_env_var(self) -> None:
        # Browser-first: a fresh browser read wins over a possibly-stale
        # env value. Users almost always want this — env var was likely
        # set ages ago when the browser path didn't work, but if browser
        # is working *now* we should prefer the live value.
        os.environ["NTULEARN_COOKIE"] = "from-env"
        with mock.patch.object(server, "read_bbrouter_cookie", return_value="from-browser"):
            self.assertEqual(server._resolve_cookie(), "from-browser")

    def test_falls_back_to_env_var_when_browser_returns_none(self) -> None:
        # Windows + Chrome/Edge ABE, browser not logged in, etc. Env var
        # is the manual safety net.
        os.environ["NTULEARN_COOKIE"] = "from-env"
        with mock.patch.object(server, "read_bbrouter_cookie", return_value=None):
            self.assertEqual(server._resolve_cookie(), "from-env")

    def test_uses_browser_when_env_unset(self) -> None:
        os.environ.pop("NTULEARN_COOKIE", None)
        with mock.patch.object(server, "read_bbrouter_cookie", return_value="from-browser"):
            self.assertEqual(server._resolve_cookie(), "from-browser")

    def test_skips_blank_env_var(self) -> None:
        # Whitespace-only NTULEARN_COOKIE is treated as unset — falls
        # through to the cache layer rather than failing validation.
        os.environ["NTULEARN_COOKIE"] = "   "
        with mock.patch.object(server, "read_bbrouter_cookie", return_value=None):
            with mock.patch.object(server, "read_cached_cookie", return_value="from-cache"):
                self.assertEqual(server._resolve_cookie(), "from-cache")

    def test_raises_when_all_sources_empty(self) -> None:
        os.environ.pop("NTULEARN_COOKIE", None)
        with mock.patch.object(server, "read_bbrouter_cookie", return_value=None):
            with self.assertRaises(RuntimeError) as ctx:
                server._resolve_cookie()
        self.assertIn("No NTULearn cookie found", str(ctx.exception))

    # --- Cache integration -----------------------------------------------
    # The hot path: every successful browser read mirrors the value into the
    # OS keychain so the next resolution can ride through a transient
    # browser-cookie3 failure (SQLite race, keychain timeout, ABE flake).

    def test_successful_browser_read_writes_value_to_cache(self) -> None:
        os.environ.pop("NTULEARN_COOKIE", None)
        with mock.patch.object(server, "read_bbrouter_cookie", return_value="fresh-value"):
            with mock.patch.object(server, "write_cached_cookie") as write:
                server._resolve_cookie()
        write.assert_called_once_with("fresh-value")

    def test_falls_back_to_cache_when_browser_and_env_both_fail(self) -> None:
        # Last-resort path: browser can't read AND no manual override is
        # set. Cache holds whatever the most-recent successful browser
        # read gave us — typically still valid (BbRouter lasts days–weeks).
        os.environ.pop("NTULEARN_COOKIE", None)
        with mock.patch.object(server, "read_bbrouter_cookie", return_value=None):
            with mock.patch.object(
                server, "read_cached_cookie", return_value="cached-value"
            ):
                self.assertEqual(server._resolve_cookie(), "cached-value")

    def test_browser_value_preferred_over_cache_when_both_present(self) -> None:
        # Cookie may have rotated since the cache was last written; the
        # browser is the source of truth, cache is just a fallback.
        os.environ.pop("NTULEARN_COOKIE", None)
        with mock.patch.object(
            server, "read_bbrouter_cookie", return_value="rotated-value"
        ):
            with mock.patch.object(
                server, "read_cached_cookie", return_value="stale-cached-value"
            ):
                self.assertEqual(server._resolve_cookie(), "rotated-value")

    def test_env_var_preferred_over_cache_when_browser_fails(self) -> None:
        # Env var is user-deliberate; cache is automatic best-guess. When
        # both are present and browser fails, env wins because it's the
        # closer-to-user-intent signal.
        os.environ["NTULEARN_COOKIE"] = "from-env"
        with mock.patch.object(server, "read_bbrouter_cookie", return_value=None):
            with mock.patch.object(
                server, "read_cached_cookie", return_value="from-cache"
            ):
                self.assertEqual(server._resolve_cookie(), "from-env")

    def test_env_var_fallback_does_not_write_cache(self) -> None:
        # We only mirror browser-derived values into the keychain. An env
        # var value might be a one-off debug cookie or a permanent manual
        # fallback — either way, persisting it isn't our job.
        os.environ["NTULEARN_COOKIE"] = "from-env"
        with mock.patch.object(server, "read_bbrouter_cookie", return_value=None):
            with mock.patch.object(server, "write_cached_cookie") as write:
                self.assertEqual(server._resolve_cookie(), "from-env")
        write.assert_not_called()

    def test_env_var_fallback_does_not_read_cache(self) -> None:
        # Env var beats cache. Once env var matches, we shouldn't even
        # consult the keychain.
        os.environ["NTULEARN_COOKIE"] = "from-env"
        with mock.patch.object(server, "read_bbrouter_cookie", return_value=None):
            with mock.patch.object(server, "read_cached_cookie") as read_cache:
                server._resolve_cookie()
        read_cache.assert_not_called()

    def test_does_not_consult_env_or_cache_when_browser_succeeds(self) -> None:
        # Hot-path optimization: when the browser read works (the common
        # case), we don't bother with env var or cache lookups.
        os.environ["NTULEARN_COOKIE"] = "from-env"
        with mock.patch.object(server, "read_bbrouter_cookie", return_value="from-browser"):
            with mock.patch.object(server, "read_cached_cookie") as read_cache:
                server._resolve_cookie()
        read_cache.assert_not_called()

    def test_raises_when_browser_env_and_cache_all_empty(self) -> None:
        os.environ.pop("NTULEARN_COOKIE", None)
        with mock.patch.object(server, "read_bbrouter_cookie", return_value=None):
            with mock.patch.object(server, "read_cached_cookie", return_value=None):
                with self.assertRaises(RuntimeError) as ctx:
                    server._resolve_cookie()
        self.assertIn("No NTULearn cookie found", str(ctx.exception))


class CookieRefreshTests(_CookieEnvIsolation, unittest.IsolatedAsyncioTestCase):
    async def asyncTearDown(self) -> None:
        # _refresh_client may have built a real httpx-backed client; close it
        # so the test's event loop doesn't complain about leftover sockets.
        if server._client is not None and server._client is not self._old_client:
            await server._client.close()

    async def test_call_tool_refreshes_then_retries_after_401(self) -> None:
        os.environ["NTULEARN_COOKIE"] = "test-cookie"
        server._client = NTULearnClient(server.BASE_URL, "test-cookie")

        call_count = 0
        original = server._list_courses

        async def flaky(client: Any, args: dict[str, Any]) -> list[TextContent]:
            nonlocal call_count
            call_count += 1
            if call_count == 1:
                raise BbRouterExpiredError()
            return [TextContent(type="text", text="success")]

        try:
            server._list_courses = flaky  # type: ignore[assignment]
            result = await server.call_tool("ntulearn_list_courses", {})
        finally:
            server._list_courses = original  # type: ignore[assignment]

        self.assertEqual(call_count, 2)
        self.assertEqual(result[0].text, "success")

    async def test_call_tool_surfaces_error_when_refresh_finds_no_cookie(self) -> None:
        os.environ["NTULEARN_COOKIE"] = "test-cookie"
        server._client = NTULearnClient(server.BASE_URL, "test-cookie")

        original = server._list_courses

        async def always_expired(client: Any, args: dict[str, Any]) -> list[TextContent]:
            raise BbRouterExpiredError()

        try:
            server._list_courses = always_expired  # type: ignore[assignment]
            # Refresh will look for a cookie and find none.
            os.environ.pop("NTULEARN_COOKIE", None)
            with mock.patch.object(server, "read_bbrouter_cookie", return_value=None):
                with self.assertRaises(RuntimeError) as ctx:
                    await server.call_tool("ntulearn_list_courses", {})
        finally:
            server._list_courses = original  # type: ignore[assignment]

        self.assertIn("No NTULearn cookie found", str(ctx.exception))

    async def test_call_tool_retries_only_once(self) -> None:
        os.environ["NTULEARN_COOKIE"] = "test-cookie"
        server._client = NTULearnClient(server.BASE_URL, "test-cookie")

        call_count = 0
        original = server._list_courses

        async def always_expired(client: Any, args: dict[str, Any]) -> list[TextContent]:
            nonlocal call_count
            call_count += 1
            raise BbRouterExpiredError()

        try:
            server._list_courses = always_expired  # type: ignore[assignment]
            with self.assertRaises(BbRouterExpiredError):
                await server.call_tool("ntulearn_list_courses", {})
        finally:
            server._list_courses = original  # type: ignore[assignment]

        # initial call + one retry, no third attempt before the error propagates
        self.assertEqual(call_count, 2)

    async def test_refresh_invalidates_cache_before_resolving(self) -> None:
        # The cookie that just produced a 401 might be the one in the cache
        # (we may have used it as a fallback when the browser failed). If we
        # re-resolve without nuking the cache first and the browser is still
        # failing, we'd loop forever on the same dead value.
        os.environ["NTULEARN_COOKIE"] = "test-cookie"
        server._client = NTULearnClient(server.BASE_URL, "test-cookie")

        with mock.patch.object(server, "delete_cached_cookie") as delete:
            await server._refresh_client()

        delete.assert_called_once()


class FakeGradebookClient:
    async def get_my_enrollments(self) -> list[dict[str, Any]]:
        return [{"courseId": "course-1", "availability": {"available": "Yes"}}]

    async def get_gradebook_columns(self, course_id: str) -> list[dict[str, Any]]:
        return [{"id": "col-1", "name": "Quiz", "score": {"possible": 10}}]

    async def get_my_user_id(self) -> str:
        return "user-1"

    async def get_user_grades(self, course_id: str, user_id: str) -> list[dict[str, Any]]:
        raise RuntimeError("grades endpoint unavailable")


class FakeDownloadClient:
    async def get_content_item(self, course_id: str, content_id: str) -> dict[str, Any]:
        return {
            "title": "Lecture",
            "contentHandler": {"id": "resource/x-bb-file"},
        }

    async def get_attachments(self, course_id: str, content_id: str) -> list[dict[str, Any]]:
        return [{"id": "att-1", "fileName": "slides.pdf"}]

    async def get_attachment_download_url(
        self, course_id: str, content_id: str, attachment_id: str
    ) -> str:
        return "/download/slides.pdf"

    async def download_bytes(self, url: str) -> tuple[bytes, str | None]:
        return b"new slides", "application/pdf"


class FakeSearchClient:
    async def get_course_contents(self, course_id: str) -> list[dict[str, Any]]:
        return [
            {"id": "1", "title": "match one", "description": {}, "hasChildren": False},
            {"id": "2", "title": "match two", "description": {}, "hasChildren": False},
            {"id": "3", "title": "match three", "description": {}, "hasChildren": False},
        ]

    async def get_content_children(
        self, course_id: str, content_id: str
    ) -> list[dict[str, Any]]:
        return []


class ToolBehaviorTests(unittest.IsolatedAsyncioTestCase):
    async def test_gradebook_partial_failure_is_reported(self) -> None:
        _, payload = await server._get_gradebook(FakeGradebookClient(), {"course_id": "course-1"})

        self.assertFalse(payload["gradesAvailable"])
        self.assertIn("grades endpoint unavailable", payload["gradeFetchError"])
        self.assertEqual(payload["columns"][0]["name"], "Quiz")
        self.assertIsNone(payload["columns"][0]["score"])

    async def test_download_does_not_overwrite_existing_file(self) -> None:
        old_download_dir = server.DOWNLOAD_DIR
        with tempfile.TemporaryDirectory() as tmp:
            server.DOWNLOAD_DIR = Path(tmp)
            (server.DOWNLOAD_DIR / "slides.pdf").write_bytes(b"old slides")
            try:
                _, payload = await server._download_file(
                    FakeDownloadClient(), {"course_id": "course-1", "content_id": "content-1"}
                )
            finally:
                server.DOWNLOAD_DIR = old_download_dir

            saved = payload["files"][0]
            self.assertEqual(saved["filename"], "slides (2).pdf")
            self.assertEqual((Path(tmp) / "slides.pdf").read_bytes(), b"old slides")
            self.assertEqual((Path(tmp) / "slides (2).pdf").read_bytes(), b"new slides")

    async def test_blank_search_query_errors(self) -> None:
        with self.assertRaises(ValueError) as ctx:
            await server._search_course_content(
                FakeSearchClient(), {"course_id": "course-1", "query": "   "}
            )
        self.assertIn("query cannot be blank", str(ctx.exception))

    async def test_search_max_results_caps_output(self) -> None:
        _, payload = await server._search_course_content(
            FakeSearchClient(),
            {"course_id": "course-1", "query": "match", "max_results": 2},
        )

        self.assertEqual(payload["count"], 2)
        self.assertEqual([item["id"] for item in payload["matches"]], ["1", "2"])


class ContentClassificationTests(unittest.TestCase):
    """Pure-function tests for _classify_kind — extension wins over content-type."""

    def test_pdf_extension_classified_as_pdf(self) -> None:
        self.assertEqual(server._classify_kind("slides.pdf", None), "pdf")

    def test_pdf_mime_classified_as_pdf(self) -> None:
        self.assertEqual(server._classify_kind("file", "application/pdf"), "pdf")

    def test_pdf_extension_wins_over_octet_stream(self) -> None:
        # The bbcswebdav case: server says octet-stream but extension is .pdf.
        self.assertEqual(
            server._classify_kind("lecture.pdf", "application/octet-stream"),
            "pdf",
        )

    def test_text_mime_classified_as_text(self) -> None:
        self.assertEqual(server._classify_kind("a", "text/plain"), "text")
        self.assertEqual(server._classify_kind("a", "text/html; charset=utf-8"), "text")
        self.assertEqual(server._classify_kind("a", "application/json"), "text")

    def test_text_extension_classified_as_text(self) -> None:
        for fname in ("notes.md", "data.csv", "config.yaml", "main.py", "readme.txt"):
            self.assertEqual(server._classify_kind(fname, None), "text", fname)

    def test_image_classified_as_binary(self) -> None:
        self.assertEqual(server._classify_kind("logo.png", "image/png"), "binary")

    def test_unknown_classified_as_binary(self) -> None:
        self.assertEqual(server._classify_kind("blob", None), "binary")
        self.assertEqual(
            server._classify_kind("video.mp4", "video/mp4"), "binary"
        )

    def test_docx_extension_classified(self) -> None:
        self.assertEqual(server._classify_kind("brief.docx", None), "docx")

    def test_pptx_extension_classified(self) -> None:
        self.assertEqual(server._classify_kind("slides.pptx", None), "pptx")

    def test_xlsx_extension_classified(self) -> None:
        self.assertEqual(server._classify_kind("data.xlsx", None), "xlsx")

    def test_office_mime_classified(self) -> None:
        cases = [
            (
                "application/vnd.openxmlformats-officedocument."
                "wordprocessingml.document",
                "docx",
            ),
            (
                "application/vnd.openxmlformats-officedocument."
                "presentationml.presentation",
                "pptx",
            ),
            (
                "application/vnd.openxmlformats-officedocument."
                "spreadsheetml.sheet",
                "xlsx",
            ),
        ]
        for mime, expected in cases:
            with self.subTest(mime=mime):
                self.assertEqual(server._classify_kind("file", mime), expected)

    def test_docx_extension_wins_over_octet_stream(self) -> None:
        # bbcswebdav case for Office files — same logic as the PDF test.
        self.assertEqual(
            server._classify_kind("brief.docx", "application/octet-stream"),
            "docx",
        )


class ContentExtractionTests(unittest.TestCase):
    """Pure-function tests for _extract_content."""

    def test_extracts_pdf_text(self) -> None:
        result = server._extract_content("hello.pdf", "application/pdf", _TINY_PDF_BYTES)
        self.assertEqual(result["kind"], "pdf")
        self.assertIn("Hello PDF Fixture", result["text"])
        self.assertEqual(result["pageCount"], 1)
        self.assertEqual(result["filename"], "hello.pdf")
        self.assertNotIn("warning", result)

    def test_decodes_utf8_text(self) -> None:
        result = server._extract_content("note.txt", "text/plain", "héllo".encode("utf-8"))
        self.assertEqual(result["kind"], "text")
        self.assertEqual(result["text"], "héllo")

    def test_decodes_with_charset_from_content_type(self) -> None:
        result = server._extract_content(
            "note.txt", "text/plain; charset=latin-1", "café".encode("latin-1")
        )
        self.assertEqual(result["text"], "café")

    def test_strips_html_tags(self) -> None:
        html = b"<html><body><h1>Hi</h1><p>Para <b>bold</b></p><script>x</script></body></html>"
        result = server._extract_content("page.html", "text/html", html)
        self.assertEqual(result["kind"], "text")
        self.assertNotIn("<h1>", result["text"])
        self.assertNotIn("<p>", result["text"])
        self.assertIn("Hi", result["text"])
        self.assertIn("bold", result["text"])

    def test_binary_returns_error_payload(self) -> None:
        png_magic = b"\x89PNG\r\n\x1a\n" + b"\x00" * 100
        result = server._extract_content("logo.png", "image/png", png_magic)
        self.assertEqual(result["kind"], "binary")
        self.assertIn("download_file", result["error"])
        self.assertIn("image/png", result["error"])

    def test_pdf_extension_wins_over_octet_stream(self) -> None:
        result = server._extract_content(
            "lecture.pdf", "application/octet-stream", _TINY_PDF_BYTES
        )
        self.assertEqual(result["kind"], "pdf")
        self.assertIn("Hello PDF Fixture", result["text"])

    def test_corrupted_pdf_returns_error(self) -> None:
        result = server._extract_content("bad.pdf", "application/pdf", b"not really a pdf")
        self.assertEqual(result["kind"], "pdf")
        self.assertIn("error", result)

    def test_extracts_docx_paragraphs(self) -> None:
        bytes_ = _make_docx_bytes(["First paragraph.", "Second paragraph."])
        result = server._extract_content("brief.docx", None, bytes_)
        self.assertEqual(result["kind"], "docx")
        self.assertIn("First paragraph", result["text"])
        self.assertIn("Second paragraph", result["text"])
        self.assertGreaterEqual(result["paragraphCount"], 2)
        self.assertEqual(result["tableCount"], 0)

    def test_extracts_docx_table(self) -> None:
        bytes_ = _make_docx_bytes(
            ["Header paragraph"],
            table_rows=[["A1", "B1"], ["A2", "B2"]],
        )
        result = server._extract_content("data.docx", None, bytes_)
        self.assertEqual(result["kind"], "docx")
        for cell in ("A1", "B1", "A2", "B2"):
            self.assertIn(cell, result["text"])
        self.assertEqual(result["tableCount"], 1)

    def test_extracts_pptx_slides(self) -> None:
        bytes_ = _make_pptx_bytes([
            {"title": "Slide One", "body": "Content one"},
            {"title": "Slide Two", "body": "Content two"},
            {"title": "Slide Three", "body": "Content three"},
        ])
        result = server._extract_content("deck.pptx", None, bytes_)
        self.assertEqual(result["kind"], "pptx")
        self.assertEqual(result["slideCount"], 3)
        self.assertIn("## Slide 1", result["text"])
        self.assertIn("## Slide 3", result["text"])
        self.assertIn("Slide One", result["text"])
        self.assertIn("Content three", result["text"])

    def test_extracts_pptx_speaker_notes(self) -> None:
        bytes_ = _make_pptx_bytes([
            {"title": "Title", "body": "Body", "notes": "Remember to mention X."},
        ])
        result = server._extract_content("deck.pptx", None, bytes_)
        self.assertEqual(result["kind"], "pptx")
        self.assertIn("Speaker notes", result["text"])
        self.assertIn("Remember to mention X", result["text"])

    def test_extracts_xlsx_sheets(self) -> None:
        bytes_ = _make_xlsx_bytes({
            "Grades": [["Student", "Score"], ["Alice", 85], ["Bob", 92]],
            "Summary": [["Mean", 88.5]],
        })
        result = server._extract_content("grades.xlsx", None, bytes_)
        self.assertEqual(result["kind"], "xlsx")
        self.assertEqual(result["sheetCount"], 2)
        self.assertIn("## Sheet: Grades", result["text"])
        self.assertIn("## Sheet: Summary", result["text"])
        self.assertIn("Alice", result["text"])
        self.assertIn("85", result["text"])
        self.assertIn("88.5", result["text"])
        self.assertNotIn("warning", result)

    def test_xlsx_truncates_above_row_cap(self) -> None:
        # One sheet exceeds the cap; assert the warning + truncation marker.
        rows = [[f"r{i}", i] for i in range(server._MAX_XLSX_ROWS_PER_SHEET + 50)]
        bytes_ = _make_xlsx_bytes({"Big": rows})
        result = server._extract_content("big.xlsx", None, bytes_)
        self.assertEqual(result["kind"], "xlsx")
        self.assertIn("warning", result)
        self.assertIn("truncated", result["text"])

    def test_corrupted_docx_returns_error(self) -> None:
        result = server._extract_content("bad.docx", None, b"not a docx file at all")
        self.assertEqual(result["kind"], "docx")
        self.assertIn("error", result)
        self.assertNotIn("text", result)


class _StubResolveClient:
    """Implements only get_content_item — for handlers that hit _resolve_content_files
    and short-circuit on no-pairs. Avoids needing get_attachments etc."""

    def __init__(self, item: dict[str, Any]) -> None:
        self._item = item

    async def get_content_item(self, course_id: str, content_id: str) -> dict[str, Any]:
        return self._item


class FakeReadClient:
    """Mock for _read_file_content that maps URLs to (bytes, content_type)."""

    def __init__(
        self,
        url_to_payload: dict[str, tuple[bytes, str | None]],
        item: dict[str, Any] | None = None,
        attachments: list[dict[str, Any]] | None = None,
        attachment_urls: dict[str, str] | None = None,
    ) -> None:
        self._url_to_payload = url_to_payload
        self._item = item or {
            "title": "Lecture 1",
            "contentHandler": {"id": "resource/x-bb-file"},
        }
        self._attachments = attachments or []
        self._attachment_urls = attachment_urls or {}

    async def get_content_item(self, course_id: str, content_id: str) -> dict[str, Any]:
        return self._item

    async def get_attachments(self, course_id: str, content_id: str) -> list[dict[str, Any]]:
        return self._attachments

    async def get_attachment_download_url(
        self, course_id: str, content_id: str, attachment_id: str
    ) -> str:
        return self._attachment_urls[attachment_id]

    async def download_bytes(self, url: str) -> tuple[bytes, str | None]:
        if url not in self._url_to_payload:
            raise AssertionError(f"unexpected download URL: {url}")
        return self._url_to_payload[url]


class ReadFileContentTests(unittest.IsolatedAsyncioTestCase):
    """End-to-end tests for the _read_file_content handler."""

    async def test_extracts_pdf_text(self) -> None:
        client = FakeReadClient(
            url_to_payload={"/url/slides.pdf": (_TINY_PDF_BYTES, "application/pdf")},
            attachments=[{"id": "att-1", "fileName": "slides.pdf"}],
            attachment_urls={"att-1": "/url/slides.pdf"},
        )
        _, payload = await server._read_file_content(
            client, {"course_id": "c", "content_id": "i"}
        )
        self.assertEqual(len(payload["files"]), 1)
        self.assertEqual(len(payload["skipped"]), 0)
        f = payload["files"][0]
        self.assertEqual(f["kind"], "pdf")
        self.assertEqual(f["filename"], "slides.pdf")
        self.assertIn("Hello PDF Fixture", f["text"])
        self.assertEqual(f["pageCount"], 1)

    async def test_decodes_text_file(self) -> None:
        client = FakeReadClient(
            url_to_payload={"/u/notes.txt": (b"hello world", "text/plain")},
            attachments=[{"id": "a", "fileName": "notes.txt"}],
            attachment_urls={"a": "/u/notes.txt"},
        )
        _, payload = await server._read_file_content(
            client, {"course_id": "c", "content_id": "i"}
        )
        self.assertEqual(payload["files"][0]["kind"], "text")
        self.assertEqual(payload["files"][0]["text"], "hello world")

    async def test_strips_html(self) -> None:
        html = b"<html><body><h1>Title</h1><p>Body text</p></body></html>"
        client = FakeReadClient(
            url_to_payload={"/u/page.html": (html, "text/html; charset=utf-8")},
            attachments=[{"id": "a", "fileName": "page.html"}],
            attachment_urls={"a": "/u/page.html"},
        )
        _, payload = await server._read_file_content(
            client, {"course_id": "c", "content_id": "i"}
        )
        text = payload["files"][0]["text"]
        self.assertNotIn("<h1>", text)
        self.assertNotIn("<body>", text)
        self.assertIn("Title", text)
        self.assertIn("Body text", text)

    async def test_refuses_binary(self) -> None:
        client = FakeReadClient(
            url_to_payload={"/u/img.png": (b"\x89PNG\r\n\x1a\n" + b"\x00" * 50, "image/png")},
            attachments=[{"id": "a", "fileName": "img.png"}],
            attachment_urls={"a": "/u/img.png"},
        )
        _, payload = await server._read_file_content(
            client, {"course_id": "c", "content_id": "i"}
        )
        self.assertEqual(payload["files"], [])
        self.assertEqual(len(payload["skipped"]), 1)
        skipped = payload["skipped"][0]
        self.assertEqual(skipped["filename"], "img.png")
        self.assertIn("download_file", skipped["reason"])

    async def test_size_cap_per_file(self) -> None:
        oversized = b"\x00" * (server._MAX_FILE_BYTES + 1)
        client = FakeReadClient(
            url_to_payload={"/u/huge.pdf": (oversized, "application/pdf")},
            attachments=[{"id": "a", "fileName": "huge.pdf"}],
            attachment_urls={"a": "/u/huge.pdf"},
        )
        _, payload = await server._read_file_content(
            client, {"course_id": "c", "content_id": "i"}
        )
        self.assertEqual(payload["files"], [])
        self.assertEqual(len(payload["skipped"]), 1)
        self.assertIn("too large", payload["skipped"][0]["reason"].lower())

    async def test_octet_stream_with_pdf_extension(self) -> None:
        # Common bbcswebdav case — server says octet-stream, extension says PDF.
        client = FakeReadClient(
            url_to_payload={
                "/u/lecture.pdf": (_TINY_PDF_BYTES, "application/octet-stream")
            },
            attachments=[{"id": "a", "fileName": "lecture.pdf"}],
            attachment_urls={"a": "/u/lecture.pdf"},
        )
        _, payload = await server._read_file_content(
            client, {"course_id": "c", "content_id": "i"}
        )
        self.assertEqual(payload["files"][0]["kind"], "pdf")
        self.assertIn("Hello PDF Fixture", payload["files"][0]["text"])

    async def test_no_files_returns_error_payload(self) -> None:
        # resource/x-bb-document with no extractable URLs in body.
        client = FakeReadClient(
            url_to_payload={},
            item={
                "title": "Empty Item",
                "contentHandler": {"id": "resource/x-bb-document"},
                "body": "<p>just some text, no links</p>",
            },
        )
        _, payload = await server._read_file_content(
            client, {"course_id": "c", "content_id": "i"}
        )
        self.assertIn("error", payload)
        self.assertEqual(payload["title"], "Empty Item")
        self.assertEqual(payload["contentHandlerId"], "resource/x-bb-document")

    async def test_per_url_401_retry(self) -> None:
        # First download raises BbRouterExpiredError; the inline retry path
        # calls _refresh_client (patched to return a fresh stub) and re-downloads.
        expired_once = {"done": False}

        class FlakyClient(FakeReadClient):
            async def download_bytes(
                self, url: str
            ) -> tuple[bytes, str | None]:
                if not expired_once["done"]:
                    expired_once["done"] = True
                    raise BbRouterExpiredError()
                return await super().download_bytes(url)

        flaky = FlakyClient(
            url_to_payload={"/u/notes.txt": (b"after retry", "text/plain")},
            attachments=[{"id": "a", "fileName": "notes.txt"}],
            attachment_urls={"a": "/u/notes.txt"},
        )

        refreshed = FakeReadClient(
            url_to_payload={"/u/notes.txt": (b"after retry", "text/plain")},
        )

        async def fake_refresh() -> Any:
            return refreshed

        with mock.patch.object(server, "_refresh_client", fake_refresh):
            _, payload = await server._read_file_content(
                flaky, {"course_id": "c", "content_id": "i"}
            )

        self.assertTrue(expired_once["done"])
        self.assertEqual(payload["files"][0]["text"], "after retry")

    async def test_handles_multiple_files(self) -> None:
        client = FakeReadClient(
            url_to_payload={
                "/u/a.txt": (b"alpha", "text/plain"),
                "/u/b.png": (b"\x89PNG\r\n\x1a\n", "image/png"),
                "/u/c.pdf": (_TINY_PDF_BYTES, "application/pdf"),
            },
            attachments=[
                {"id": "1", "fileName": "a.txt"},
                {"id": "2", "fileName": "b.png"},
                {"id": "3", "fileName": "c.pdf"},
            ],
            attachment_urls={"1": "/u/a.txt", "2": "/u/b.png", "3": "/u/c.pdf"},
        )
        _, payload = await server._read_file_content(
            client, {"course_id": "c", "content_id": "i"}
        )
        kinds = sorted(f["kind"] for f in payload["files"])
        self.assertEqual(kinds, ["pdf", "text"])
        self.assertEqual(len(payload["skipped"]), 1)
        self.assertEqual(payload["skipped"][0]["filename"], "b.png")

    async def test_reads_docx_attachment(self) -> None:
        # End-to-end happy path for one Office format. The per-format unit
        # tests in ContentExtractionTests cover pptx/xlsx parsing details;
        # here we just verify the handler routes Office files into `files`
        # (not `skipped`) and threads through the content_type override.
        docx_bytes = _make_docx_bytes(["Assignment brief content."])
        client = FakeReadClient(
            url_to_payload={
                "/u/brief.docx": (docx_bytes, "application/octet-stream")
            },
            attachments=[{"id": "a", "fileName": "brief.docx"}],
            attachment_urls={"a": "/u/brief.docx"},
        )
        _, payload = await server._read_file_content(
            client, {"course_id": "c", "content_id": "i"}
        )
        self.assertEqual(len(payload["files"]), 1)
        self.assertEqual(payload["skipped"], [])
        f = payload["files"][0]
        self.assertEqual(f["kind"], "docx")
        self.assertIn("Assignment brief content", f["text"])


class ResolveContentFilesTests(unittest.IsolatedAsyncioTestCase):
    """Direct tests of the _resolve_content_files helper."""

    async def test_resolves_attachment_handler(self) -> None:
        client = FakeReadClient(
            url_to_payload={},
            attachments=[
                {"id": "a", "fileName": "x.pdf"},
                {"id": "b", "fileName": "y.docx"},
            ],
            attachment_urls={"a": "/u/x.pdf", "b": "/u/y.docx"},
        )
        item, handler_id, pairs = await server._resolve_content_files(client, "c", "i")
        self.assertEqual(handler_id, "resource/x-bb-file")
        self.assertEqual(pairs, [("/u/x.pdf", "x.pdf"), ("/u/y.docx", "y.docx")])
        self.assertEqual(item["title"], "Lecture 1")

    async def test_resolves_html_body_handler(self) -> None:
        # Use the real extract_all_files path via a body containing bbcswebdav links.
        body = (
            '<a href="https://ntulearn.ntu.edu.sg/bbcswebdav/'
            'pid-1/notes.pdf">Notes</a>'
        )
        client = FakeReadClient(
            url_to_payload={},
            item={
                "title": "Doc",
                "contentHandler": {"id": "resource/x-bb-document"},
                "body": body,
            },
        )
        _, handler_id, pairs = await server._resolve_content_files(client, "c", "i")
        self.assertEqual(handler_id, "resource/x-bb-document")
        self.assertEqual(len(pairs), 1)
        self.assertIn("bbcswebdav", pairs[0][0])

    async def test_no_files_returns_empty_pairs(self) -> None:
        client = FakeReadClient(
            url_to_payload={},
            item={
                "title": "Empty",
                "contentHandler": {"id": "resource/x-bb-document"},
                "body": "<p>nothing useful</p>",
            },
        )
        _, _, pairs = await server._resolve_content_files(client, "c", "i")
        self.assertEqual(pairs, [])


class PdfModeAndPageRangeParsingTests(unittest.TestCase):
    """Pure-function tests for the PDF mode/pages arg parsers."""

    def test_resolve_pdf_mode_default_is_text(self) -> None:
        self.assertEqual(server._resolve_pdf_mode({}), "text")

    def test_resolve_pdf_mode_explicit_text(self) -> None:
        self.assertEqual(server._resolve_pdf_mode({"mode": "text"}), "text")

    def test_resolve_pdf_mode_explicit_vision(self) -> None:
        self.assertEqual(server._resolve_pdf_mode({"mode": "vision"}), "vision")

    def test_resolve_pdf_mode_auto_is_alias_for_text(self) -> None:
        # 'auto' kept for backwards compatibility with earlier callers.
        self.assertEqual(server._resolve_pdf_mode({"mode": "auto"}), "text")
        self.assertEqual(server._resolve_pdf_mode({"mode": "AUTO"}), "text")

    def test_resolve_pdf_mode_invalid_raises(self) -> None:
        with self.assertRaises(ValueError):
            server._resolve_pdf_mode({"mode": "garbage"})

    def test_parse_page_range_none(self) -> None:
        self.assertIsNone(server._parse_page_range(None))

    def test_parse_page_range_empty(self) -> None:
        self.assertIsNone(server._parse_page_range(""))
        self.assertIsNone(server._parse_page_range("   "))

    def test_parse_page_range_single_page(self) -> None:
        self.assertEqual(server._parse_page_range("5"), {5})

    def test_parse_page_range_simple_range(self) -> None:
        self.assertEqual(server._parse_page_range("1-3"), {1, 2, 3})

    def test_parse_page_range_mixed(self) -> None:
        self.assertEqual(
            server._parse_page_range("1-3,5,8-9"), {1, 2, 3, 5, 8, 9}
        )

    def test_parse_page_range_whitespace_tolerated(self) -> None:
        self.assertEqual(server._parse_page_range(" 1 - 3 , 5 "), {1, 2, 3, 5})

    def test_parse_page_range_invalid_token_raises(self) -> None:
        with self.assertRaises(ValueError):
            server._parse_page_range("abc")

    def test_parse_page_range_zero_raises(self) -> None:
        with self.assertRaises(ValueError):
            server._parse_page_range("0")

    def test_parse_page_range_inverted_raises(self) -> None:
        with self.assertRaises(ValueError):
            server._parse_page_range("5-3")


class PdfVisionExtractionTests(unittest.TestCase):
    """Pure-function tests for _extract_pdf_vision."""

    def test_renders_each_page_as_png(self) -> None:
        result = server._extract_pdf_vision(
            "slides.pdf", "application/pdf", _TINY_PDF_BYTES, len(_TINY_PDF_BYTES), None
        )
        self.assertEqual(result["kind"], "pdf")
        self.assertEqual(result["pageCount"], 1)
        self.assertEqual(result["pagesRendered"], [1])
        images = result["_images"]
        self.assertEqual(len(images), 1)
        label, png_bytes = images[0]
        self.assertIn("page 1", label)
        # PNG magic header — confirms PyMuPDF actually rendered an image.
        self.assertTrue(png_bytes.startswith(b"\x89PNG\r\n\x1a\n"))

    def test_text_includes_per_page_header(self) -> None:
        result = server._extract_pdf_vision(
            "slides.pdf", "application/pdf", _TINY_PDF_BYTES, len(_TINY_PDF_BYTES), None
        )
        self.assertIn("## Page 1", result["text"])
        self.assertIn("Hello PDF Fixture", result["text"])

    def test_pages_filter_restricts_render(self) -> None:
        # Filter to a non-existent page — extractor should render nothing
        # and report empty pagesRendered. This exercises the 1-indexed
        # bounds check without needing a multi-page fixture.
        result = server._extract_pdf_vision(
            "slides.pdf",
            "application/pdf",
            _TINY_PDF_BYTES,
            len(_TINY_PDF_BYTES),
            {99},
        )
        self.assertEqual(result["pageCount"], 1)
        self.assertEqual(result["pagesRendered"], [])
        self.assertEqual(result["_images"], [])

    def test_corrupted_pdf_returns_error(self) -> None:
        result = server._extract_pdf_vision(
            "bad.pdf", "application/pdf", b"not really a pdf", 16, None
        )
        self.assertEqual(result["kind"], "pdf")
        self.assertIn("error", result)
        self.assertNotIn("_images", result)


class ReadFileContentVisionTests(unittest.IsolatedAsyncioTestCase):
    """End-to-end tests for the vision/text mode handling in _read_file_content."""

    async def test_default_mode_is_text(self) -> None:
        client = FakeReadClient(
            url_to_payload={"/u/slides.pdf": (_TINY_PDF_BYTES, "application/pdf")},
            attachments=[{"id": "a", "fileName": "slides.pdf"}],
            attachment_urls={"a": "/u/slides.pdf"},
        )
        content, payload = await server._read_file_content(
            client, {"course_id": "c", "content_id": "i"}
        )
        f = payload["files"][0]
        self.assertEqual(f["kind"], "pdf")
        # Default is text — pypdf path, no pages rendered.
        self.assertNotIn("pagesRendered", f)
        self.assertIn("Hello PDF Fixture", f["text"])
        image_blocks = [b for b in content if isinstance(b, ImageContent)]
        self.assertEqual(image_blocks, [])

    async def test_vision_mode_emits_image_blocks(self) -> None:
        client = FakeReadClient(
            url_to_payload={"/u/slides.pdf": (_TINY_PDF_BYTES, "application/pdf")},
            attachments=[{"id": "a", "fileName": "slides.pdf"}],
            attachment_urls={"a": "/u/slides.pdf"},
        )
        content, payload = await server._read_file_content(
            client, {"course_id": "c", "content_id": "i", "mode": "vision"}
        )
        # Structured payload still has the text + page count.
        f = payload["files"][0]
        self.assertEqual(f["kind"], "pdf")
        self.assertEqual(f["pageCount"], 1)
        self.assertEqual(f["pagesRendered"], [1])
        # Structured payload must NOT carry the raw image bytes — those
        # only belong in the unstructured content list as ImageContent.
        self.assertNotIn("_images", f)
        # Unstructured content includes one TextContent + one ImageContent.
        text_blocks = [b for b in content if isinstance(b, TextContent)]
        image_blocks = [b for b in content if isinstance(b, ImageContent)]
        self.assertEqual(len(text_blocks), 1)
        self.assertEqual(len(image_blocks), 1)
        self.assertEqual(image_blocks[0].mimeType, "image/png")
        # base64 data is non-empty and decodes back to a PNG.
        self.assertTrue(image_blocks[0].data)
        self.assertTrue(
            base64.b64decode(image_blocks[0].data).startswith(b"\x89PNG\r\n\x1a\n")
        )

    async def test_text_mode_skips_image_rendering(self) -> None:
        client = FakeReadClient(
            url_to_payload={"/u/slides.pdf": (_TINY_PDF_BYTES, "application/pdf")},
            attachments=[{"id": "a", "fileName": "slides.pdf"}],
            attachment_urls={"a": "/u/slides.pdf"},
        )
        content, payload = await server._read_file_content(
            client, {"course_id": "c", "content_id": "i", "mode": "text"}
        )
        f = payload["files"][0]
        self.assertEqual(f["kind"], "pdf")
        # text-mode pypdf path does not populate pagesRendered.
        self.assertNotIn("pagesRendered", f)
        self.assertIn("Hello PDF Fixture", f["text"])
        # No ImageContent blocks in the unstructured content list.
        image_blocks = [b for b in content if isinstance(b, ImageContent)]
        self.assertEqual(image_blocks, [])

    async def test_invalid_mode_raises(self) -> None:
        client = FakeReadClient(
            url_to_payload={"/u/slides.pdf": (_TINY_PDF_BYTES, "application/pdf")},
            attachments=[{"id": "a", "fileName": "slides.pdf"}],
            attachment_urls={"a": "/u/slides.pdf"},
        )
        with self.assertRaises(ValueError):
            await server._read_file_content(
                client,
                {"course_id": "c", "content_id": "i", "mode": "garbage"},
            )

    async def test_pages_filter_threads_through(self) -> None:
        # Single-page PDF; pages="1" preserves it. Asserts the arg actually
        # reaches _extract_pdf_vision (vs. being silently dropped).
        client = FakeReadClient(
            url_to_payload={"/u/slides.pdf": (_TINY_PDF_BYTES, "application/pdf")},
            attachments=[{"id": "a", "fileName": "slides.pdf"}],
            attachment_urls={"a": "/u/slides.pdf"},
        )
        _, payload = await server._read_file_content(
            client,
            {"course_id": "c", "content_id": "i", "mode": "vision", "pages": "1"},
        )
        self.assertEqual(payload["files"][0]["pagesRendered"], [1])

    async def test_pages_filter_excludes_out_of_range(self) -> None:
        client = FakeReadClient(
            url_to_payload={"/u/slides.pdf": (_TINY_PDF_BYTES, "application/pdf")},
            attachments=[{"id": "a", "fileName": "slides.pdf"}],
            attachment_urls={"a": "/u/slides.pdf"},
        )
        content, payload = await server._read_file_content(
            client,
            {"course_id": "c", "content_id": "i", "mode": "vision", "pages": "99"},
        )
        # Page 99 doesn't exist on a 1-page PDF — extractor should render
        # nothing and we should see no ImageContent.
        self.assertEqual(payload["files"][0]["pagesRendered"], [])
        image_blocks = [b for b in content if isinstance(b, ImageContent)]
        self.assertEqual(image_blocks, [])

    async def test_invalid_pages_arg_raises(self) -> None:
        client = FakeReadClient(
            url_to_payload={"/u/slides.pdf": (_TINY_PDF_BYTES, "application/pdf")},
            attachments=[{"id": "a", "fileName": "slides.pdf"}],
            attachment_urls={"a": "/u/slides.pdf"},
        )
        with self.assertRaises(ValueError):
            await server._read_file_content(
                client,
                {"course_id": "c", "content_id": "i", "pages": "abc"},
            )

    async def test_text_mode_does_not_apply_to_office_formats(self) -> None:
        # mode='text' is only meaningful for PDFs; .docx still extracts
        # paragraphs via python-docx and produces no ImageContent.
        docx_bytes = _make_docx_bytes(["Sample text body."])
        client = FakeReadClient(
            url_to_payload={"/u/brief.docx": (docx_bytes, None)},
            attachments=[{"id": "a", "fileName": "brief.docx"}],
            attachment_urls={"a": "/u/brief.docx"},
        )
        content, payload = await server._read_file_content(
            client, {"course_id": "c", "content_id": "i", "mode": "text"}
        )
        self.assertEqual(payload["files"][0]["kind"], "docx")
        self.assertIn("Sample text body", payload["files"][0]["text"])
        image_blocks = [b for b in content if isinstance(b, ImageContent)]
        self.assertEqual(image_blocks, [])


class FakeCalendarClient:
    """Minimal client surface for _get_upcoming cross-course fan-out tests."""

    def __init__(
        self,
        enrollments: list[dict[str, Any]],
        per_course_items: dict[str, list[dict[str, Any]]],
        failing_courses: set[str] | None = None,
    ) -> None:
        self._enrollments = enrollments
        self._per_course = per_course_items
        self._failing = failing_courses or set()
        self.calls: list[dict[str, Any]] = []

    async def get_my_enrollments(self) -> list[dict[str, Any]]:
        return self._enrollments

    async def get_calendar_items(
        self,
        *,
        course_id: str | None = None,
        since: str | None = None,
        until: str | None = None,
        item_type: str | None = None,
    ) -> list[dict[str, Any]]:
        self.calls.append(
            {"course_id": course_id, "since": since, "until": until, "item_type": item_type}
        )
        if course_id in self._failing:
            raise RuntimeError(f"course {course_id} blew up")
        return self._per_course.get(course_id or "", [])


class UpcomingTests(unittest.IsolatedAsyncioTestCase):
    async def test_fans_out_across_enrolled_courses_by_default(self) -> None:
        client = FakeCalendarClient(
            enrollments=[
                {"courseId": "_1_1", "availability": {"available": "Yes"}},
                {"courseId": "_2_1", "availability": {"available": "Yes"}},
            ],
            per_course_items={
                "_1_1": [
                    {
                        "id": "ci-a",
                        "calendarName": "course-cal-1",
                        "title": "Lecture A",
                        "start": "2026-05-20T10:00:00Z",
                        "end": "2026-05-20T11:00:00Z",
                    }
                ],
                "_2_1": [
                    {
                        "id": "ci-b",
                        "calendarName": "course-cal-2",
                        "title": "Quiz B",
                        "start": "2026-05-18T09:00:00Z",
                        "end": "2026-05-18T10:00:00Z",
                        "dynamicCalendarItemProps": {
                            "eventType": "GradebookColumn",
                            "gradable": True,
                            "attemptable": True,
                        },
                    }
                ],
            },
        )
        _, payload = await server._get_upcoming(client, {})

        self.assertEqual(sorted(payload["courseIdsQueried"]), ["_1_1", "_2_1"])
        ids = [i["id"] for i in payload["items"]]
        self.assertEqual(ids, ["ci-b", "ci-a"])  # sorted by start asc
        self.assertEqual(payload["items"][0]["courseId"], "_2_1")
        self.assertTrue(payload["items"][0]["gradable"])
        self.assertEqual(payload["courseErrors"], {})

    async def test_explicit_course_ids_overrides_enrollment_fanout(self) -> None:
        client = FakeCalendarClient(
            enrollments=[
                {"courseId": "_1_1", "availability": {"available": "Yes"}},
                {"courseId": "_2_1", "availability": {"available": "Yes"}},
            ],
            per_course_items={"_2_1": []},
        )
        await server._get_upcoming(client, {"course_ids": ["_2_1"]})
        self.assertEqual([c["course_id"] for c in client.calls], ["_2_1"])

    async def test_type_filter_threads_through(self) -> None:
        client = FakeCalendarClient(
            enrollments=[{"courseId": "_1_1", "availability": {"available": "Yes"}}],
            per_course_items={"_1_1": []},
        )
        await server._get_upcoming(
            client,
            {"course_ids": ["_1_1"], "type": "GradebookColumn"},
        )
        self.assertEqual(client.calls[0]["item_type"], "GradebookColumn")

    async def test_invalid_iso8601_since_raises(self) -> None:
        client = FakeCalendarClient(
            enrollments=[{"courseId": "_1_1", "availability": {"available": "Yes"}}],
            per_course_items={"_1_1": []},
        )
        with self.assertRaises(ValueError):
            await server._get_upcoming(client, {"since": "not-a-date"})

    async def test_per_course_failure_recorded(self) -> None:
        client = FakeCalendarClient(
            enrollments=[
                {"courseId": "_1_1", "availability": {"available": "Yes"}},
                {"courseId": "_2_1", "availability": {"available": "Yes"}},
            ],
            per_course_items={"_1_1": []},
            failing_courses={"_2_1"},
        )
        _, payload = await server._get_upcoming(client, {})
        self.assertEqual(payload["items"], [])
        self.assertIn("_2_1", payload["courseErrors"])
        self.assertIn("blew up", payload["courseErrors"]["_2_1"])


class FakeCrossCourseAnnouncementsClient:
    """Client that returns per-course announcements + enrollments."""

    def __init__(
        self,
        enrollments: list[dict[str, Any]],
        per_course: dict[str, list[dict[str, Any]]],
        failing: set[str] | None = None,
    ) -> None:
        self._enrollments = enrollments
        self._per_course = per_course
        self._failing = failing or set()

    async def get_my_enrollments(self) -> list[dict[str, Any]]:
        return self._enrollments

    async def get_announcements(self, course_id: str) -> list[dict[str, Any]]:
        if course_id in self._failing:
            raise RuntimeError(f"forbidden for {course_id}")
        return self._per_course.get(course_id, [])


class AnnouncementsCrossCourseTests(unittest.IsolatedAsyncioTestCase):
    def _ann(self, ann_id: str, created: str) -> dict[str, Any]:
        return {
            "id": ann_id,
            "title": ann_id,
            "body": {"rawText": "<p>body</p>"},
            "created": created,
            "modified": created,
            "availability": {"available": "Yes"},
        }

    async def test_default_fans_out_and_sorts_newest_first(self) -> None:
        client = FakeCrossCourseAnnouncementsClient(
            enrollments=[
                {"courseId": "_1_1", "availability": {"available": "Yes"}},
                {"courseId": "_2_1", "availability": {"available": "Yes"}},
            ],
            per_course={
                "_1_1": [self._ann("a1", "2026-05-15T10:00:00Z")],
                "_2_1": [self._ann("a2", "2026-05-16T10:00:00Z")],
            },
        )
        _, payload = await server._get_announcements(client, {})
        ids = [a["id"] for a in payload["announcements"]]
        self.assertEqual(ids, ["a2", "a1"])
        self.assertEqual(payload["announcements"][0]["courseId"], "_2_1")
        self.assertEqual(payload["announcements"][1]["courseId"], "_1_1")
        self.assertEqual(sorted(payload["courseIdsQueried"]), ["_1_1", "_2_1"])
        self.assertEqual(payload["courseErrors"], {})

    async def test_since_filters_older_announcements(self) -> None:
        client = FakeCrossCourseAnnouncementsClient(
            enrollments=[{"courseId": "_1_1", "availability": {"available": "Yes"}}],
            per_course={
                "_1_1": [
                    self._ann("old", "2025-01-01T00:00:00Z"),
                    self._ann("new", "2026-05-15T00:00:00Z"),
                ],
            },
        )
        _, payload = await server._get_announcements(
            client, {"since": "2026-05-09T00:00:00Z"}
        )
        ids = [a["id"] for a in payload["announcements"]]
        self.assertEqual(ids, ["new"])

    async def test_explicit_course_ids_scopes_to_one_course(self) -> None:
        client = FakeCrossCourseAnnouncementsClient(
            enrollments=[
                {"courseId": "_1_1", "availability": {"available": "Yes"}},
                {"courseId": "_2_1", "availability": {"available": "Yes"}},
            ],
            per_course={
                "_1_1": [self._ann("a1", "2026-05-15T00:00:00Z")],
                "_2_1": [self._ann("a2", "2026-05-15T00:00:00Z")],
            },
        )
        _, payload = await server._get_announcements(
            client, {"course_ids": ["_1_1"]}
        )
        ids = [a["id"] for a in payload["announcements"]]
        self.assertEqual(ids, ["a1"])
        self.assertEqual(payload["courseIdsQueried"], ["_1_1"])

    async def test_per_course_failure_recorded_not_raised(self) -> None:
        client = FakeCrossCourseAnnouncementsClient(
            enrollments=[
                {"courseId": "_1_1", "availability": {"available": "Yes"}},
                {"courseId": "_2_1", "availability": {"available": "Yes"}},
            ],
            per_course={"_1_1": [self._ann("a1", "2026-05-15T00:00:00Z")]},
            failing={"_2_1"},
        )
        _, payload = await server._get_announcements(client, {})
        self.assertEqual([a["id"] for a in payload["announcements"]], ["a1"])
        self.assertIn("_2_1", payload["courseErrors"])

    async def test_invalid_since_raises(self) -> None:
        client = FakeCrossCourseAnnouncementsClient(
            enrollments=[{"courseId": "_1_1", "availability": {"available": "Yes"}}],
            per_course={"_1_1": []},
        )
        with self.assertRaises(ValueError):
            await server._get_announcements(client, {"since": "garbage"})


class FakeCrossCourseGradebookClient:
    def __init__(
        self,
        enrollments: list[dict[str, Any]],
        per_course_columns: dict[str, list[dict[str, Any]]],
        per_course_grades: dict[str, list[dict[str, Any]]],
        failing_columns: set[str] | None = None,
        grades_fail_globally: bool = False,
    ) -> None:
        self._enrollments = enrollments
        self._cols = per_course_columns
        self._grades = per_course_grades
        self._failing_cols = failing_columns or set()
        self._grades_fail_globally = grades_fail_globally

    async def get_my_enrollments(self) -> list[dict[str, Any]]:
        return self._enrollments

    async def get_my_user_id(self) -> str:
        return "user-1"

    async def get_gradebook_columns(self, course_id: str) -> list[dict[str, Any]]:
        if course_id in self._failing_cols:
            raise RuntimeError(f"columns failed for {course_id}")
        return self._cols.get(course_id, [])

    async def get_user_grades(
        self, course_id: str, user_id: str
    ) -> list[dict[str, Any]]:
        if self._grades_fail_globally:
            raise RuntimeError(f"grades endpoint unavailable for {course_id}")
        return self._grades.get(course_id, [])


class GradebookCrossCourseTests(unittest.IsolatedAsyncioTestCase):
    async def test_default_fans_out_and_attributes_courseid_per_column(self) -> None:
        client = FakeCrossCourseGradebookClient(
            enrollments=[
                {"courseId": "_1_1", "availability": {"available": "Yes"}},
                {"courseId": "_2_1", "availability": {"available": "Yes"}},
            ],
            per_course_columns={
                "_1_1": [
                    {"id": "col-a", "name": "Quiz 1", "score": {"possible": 10}}
                ],
                "_2_1": [
                    {"id": "col-b", "name": "Lab 1", "score": {"possible": 20}}
                ],
            },
            per_course_grades={
                "_1_1": [{"columnId": "col-a", "score": 9, "grade": "9"}],
                "_2_1": [],
            },
        )
        _, payload = await server._get_gradebook(client, {})
        by_id = {c["id"]: c for c in payload["columns"]}
        self.assertEqual(by_id["col-a"]["courseId"], "_1_1")
        self.assertEqual(by_id["col-b"]["courseId"], "_2_1")
        self.assertEqual(by_id["col-a"]["score"], 9)
        self.assertTrue(payload["gradesAvailable"])

    async def test_column_failure_per_course_recorded(self) -> None:
        client = FakeCrossCourseGradebookClient(
            enrollments=[
                {"courseId": "_1_1", "availability": {"available": "Yes"}},
                {"courseId": "_2_1", "availability": {"available": "Yes"}},
            ],
            per_course_columns={
                "_1_1": [{"id": "col-a", "name": "Quiz", "score": {"possible": 10}}],
            },
            per_course_grades={"_1_1": []},
            failing_columns={"_2_1"},
        )
        _, payload = await server._get_gradebook(client, {})
        self.assertEqual([c["id"] for c in payload["columns"]], ["col-a"])
        self.assertIn("_2_1", payload["courseErrors"])

    async def test_global_grade_failure_flips_grades_available(self) -> None:
        client = FakeCrossCourseGradebookClient(
            enrollments=[
                {"courseId": "_1_1", "availability": {"available": "Yes"}},
            ],
            per_course_columns={
                "_1_1": [{"id": "col-a", "name": "Quiz", "score": {"possible": 10}}],
            },
            per_course_grades={},
            grades_fail_globally=True,
        )
        _, payload = await server._get_gradebook(client, {})
        self.assertFalse(payload["gradesAvailable"])
        self.assertIn("grades endpoint unavailable", payload["gradeFetchError"])
        self.assertEqual(payload["columns"][0]["score"], None)


class FakeContentsClient:
    def __init__(
        self,
        top_level: list[dict[str, Any]],
        children_by_id: dict[str, list[dict[str, Any]]] | None = None,
    ) -> None:
        self._top = top_level
        self._children = children_by_id or {}
        self.calls: list[tuple[str, str | None]] = []

    async def get_course_contents(self, course_id: str) -> list[dict[str, Any]]:
        self.calls.append((course_id, None))
        return self._top

    async def get_content_children(
        self, course_id: str, content_id: str
    ) -> list[dict[str, Any]]:
        self.calls.append((course_id, content_id))
        return self._children.get(content_id, [])


class GetCourseContentsTests(unittest.IsolatedAsyncioTestCase):
    async def test_without_parent_id_returns_root(self) -> None:
        client = FakeContentsClient(
            top_level=[{"id": "x", "title": "X"}],
        )
        _, payload = await server._get_course_contents(
            client, {"course_id": "_1_1"}
        )
        self.assertEqual(payload["items"][0]["id"], "x")
        self.assertEqual(client.calls, [("_1_1", None)])

    async def test_with_parent_id_drills_into_children(self) -> None:
        client = FakeContentsClient(
            top_level=[],
            children_by_id={
                "_p_1": [{"id": "kid-1", "title": "Kid"}],
            },
        )
        _, payload = await server._get_course_contents(
            client, {"course_id": "_1_1", "parent_id": "_p_1"}
        )
        self.assertEqual(payload["items"][0]["id"], "kid-1")
        self.assertEqual(client.calls, [("_1_1", "_p_1")])


class FakeDownloadClientForDest:
    """download-only fake that records bytes per URL."""

    def __init__(self, items: list[tuple[str, str]]) -> None:
        # items: list of (url, filename)
        self._items = items

    async def get_content_item(self, course_id: str, content_id: str) -> dict[str, Any]:
        return {
            "title": "Lecture pack",
            "contentHandler": {"id": "resource/x-bb-file"},
        }

    async def get_attachments(
        self, course_id: str, content_id: str
    ) -> list[dict[str, Any]]:
        return [{"id": f"att-{i}", "fileName": fn} for i, (_, fn) in enumerate(self._items)]

    async def get_attachment_download_url(
        self, course_id: str, content_id: str, attachment_id: str
    ) -> str:
        idx = int(attachment_id.split("-")[1])
        return self._items[idx][0]

    async def download_bytes(self, url: str) -> tuple[bytes, str | None]:
        return b"payload", "application/pdf"


class DownloadDestinationTests(unittest.IsolatedAsyncioTestCase):
    async def test_explicit_destination_dir_used(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            target = Path(tmp) / "y3s1" / "sc2002" / "week 8"
            _, payload = await server._download_file(
                FakeDownloadClientForDest([("/u/tut.pdf", "tut.pdf")]),
                {
                    "course_id": "_1_1",
                    "content_id": "_2_1",
                    "destination_dir": str(target),
                },
            )
            self.assertTrue(target.exists())
            saved_path = Path(payload["files"][0]["localPath"])
            self.assertEqual(saved_path.parent.resolve(), target.resolve())
            self.assertEqual(saved_path.read_bytes(), b"payload")

    async def test_tilde_destination_expands(self) -> None:
        # We don't want to actually write under ~ in tests; just verify the
        # resolver expansion logic.
        resolved = server._resolve_destination_dir("~/foo/bar")
        self.assertFalse(str(resolved).startswith("~"))
        self.assertTrue(str(resolved).endswith("foo/bar"))

    async def test_env_var_default_when_no_arg(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            tmp_resolved = Path(tmp).resolve()
            old = os.environ.get("NTULEARN_DOWNLOAD_DIR")
            os.environ["NTULEARN_DOWNLOAD_DIR"] = tmp
            try:
                _, payload = await server._download_file(
                    FakeDownloadClientForDest([("/u/tut.pdf", "tut.pdf")]),
                    {"course_id": "_1_1", "content_id": "_2_1"},
                )
            finally:
                if old is None:
                    del os.environ["NTULEARN_DOWNLOAD_DIR"]
                else:
                    os.environ["NTULEARN_DOWNLOAD_DIR"] = old
            saved_parent = Path(payload["files"][0]["localPath"]).resolve().parent
            self.assertEqual(saved_parent, tmp_resolved)

    async def test_empty_destination_dir_raises(self) -> None:
        with self.assertRaises(ValueError):
            server._resolve_destination_dir("   ")

    async def test_non_string_destination_dir_raises(self) -> None:
        with self.assertRaises(ValueError):
            server._resolve_destination_dir(123)


class PDFTextDefaultTests(unittest.IsolatedAsyncioTestCase):
    """Step 6 acceptance: PDFs default to text, no ImageContent emitted."""

    async def test_pdf_no_mode_arg_returns_text_only(self) -> None:
        client = FakeReadClient(
            url_to_payload={"/u/notes.pdf": (_TINY_PDF_BYTES, "application/pdf")},
            attachments=[{"id": "a", "fileName": "notes.pdf"}],
            attachment_urls={"a": "/u/notes.pdf"},
        )
        content, payload = await server._read_file_content(
            client, {"course_id": "c", "content_id": "i"}
        )
        f = payload["files"][0]
        self.assertEqual(f["kind"], "pdf")
        self.assertIn("Hello PDF Fixture", f["text"])
        self.assertEqual(
            [b for b in content if isinstance(b, ImageContent)], []
        )


class PDFStdoutProtectionTests(unittest.IsolatedAsyncioTestCase):
    """Step 1 acceptance: nothing leaks to stdout during PDF vision render.

    Captures fd=1 across the render call by redirecting it to a pipe, then
    asserts the pipe is empty. This is a stronger guarantee than mocking out
    `print` — it catches C-level emissions from MuPDF.
    """

    async def test_vision_render_does_not_write_to_stdout(self) -> None:
        client = FakeReadClient(
            url_to_payload={"/u/notes.pdf": (_TINY_PDF_BYTES, "application/pdf")},
            attachments=[{"id": "a", "fileName": "notes.pdf"}],
            attachment_urls={"a": "/u/notes.pdf"},
        )

        # Save the real stdout fd, redirect fd=1 to a pipe for the render call.
        original_stdout_fd = os.dup(1)
        r_fd, w_fd = os.pipe()
        os.dup2(w_fd, 1)
        try:
            try:
                await server._read_file_content(
                    client,
                    {"course_id": "c", "content_id": "i", "mode": "vision"},
                )
            finally:
                # Restore real stdout before reading the pipe — otherwise the
                # subsequent assert prints would also land in the pipe.
                os.dup2(original_stdout_fd, 1)
                os.close(original_stdout_fd)
                os.close(w_fd)

            captured = os.read(r_fd, 8192)
        finally:
            os.close(r_fd)

        self.assertEqual(captured, b"")


if __name__ == "__main__":
    unittest.main()
